from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.fill import FillFormat

# Создаём презентацию
prs = Presentation()

# Функция добавления слайда с заголовком и содержимым с улучшенным дизайном
def add_slide(title, content, images=None, bg_color=RGBColor(240, 248, 255)):
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
    slide = prs.slides.add_slide(slide_layout)

    # Фон слайда
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.add_paragraph()
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 25, 112)

    # Основной текст
    top = 1.5
    for paragraph in content:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(6.5), Inches(0.7))
        tf = text_box.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = paragraph
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT
        top += 0.8

    # Добавляем изображения
    if images:
        left = Inches(7.2)
        top_img = Inches(1.5)
        for img_path in images:
            slide.shapes.add_picture(img_path, left, top_img, width=Inches(2.5))
            top_img += Inches(2.2)

# Слайды с контентом и изображениями
slides_content = [
    ("Neural Style Transfer Web Application", [
        "Upload your photo and apply famous artist styles.",
        "Built with Python Flask and PyTorch.",
        "Minimalist interface with progress tracking and style preview."
    ], None),
    ("Key Features", [
        "Upload a photo and choose a style from 7 famous artists (Van Gogh, Picasso, Dali, Munch, Klimt, Chagall, Mucha).",
        "Style preview before generating.",
        "Soft style application preserves original details.",
        "Progress bar shows generation status.",
        "Single session image processing, no history stored."
    ], [
        'static/models/munch.jpg',
        'static/uploads/0f0ca4848c5a45c38e1ffdf2ded2789c_2025-10-19_22.08.11.jpg'
    ]),
    ("Architecture", [
        "Frontend: HTML + Bootstrap, upload form, style preview, progress bar, result display.",
        "Backend: Flask server, handles uploads and style transfer, SSE for progress, logs to server.log.",
        "Neural Style Transfer Engine: PyTorch, VGG19 preloaded, soft style application configurable via weights."
    ], None),
    ("Workflow", [
        "1. Upload photo.",
        "2. Select artist style.",
        "3. Frontend shows style preview.",
        "4. Click Generate, backend processes NST in a separate thread.",
        "5. Progress bar updates in real-time.",
        "6. Final styled image displayed on the same page."
    ], None),
    ("Technology Stack & Notes", [
        "Python 3, Flask, PyTorch / Torchvision, PIL (Pillow), Bootstrap 5.",
        "Simple, lightweight, and easy to deploy.",
        "Designed for single-user sessions.",
        "Style intensity adjustable via weights in NST.",
        "Can be deployed on GitHub Pages + Heroku/Docker for demo purposes."
    ], None),
    ("How NST Works", [
        "Input: Content image + Style image.",
        "VGG19 extracts content and style features.",
        "Content loss preserves structure; Style loss transfers colors/textures.",
        "Gradient-based optimization updates input image iteratively.",
        "Output: Stylized image preserving original content." 
    ], None)
]

for title, content, images in slides_content:
    add_slide(title, content, images)

# Сохраняем презентацию
prs.save('Neural_Style_Transfer_Presentation_Enhanced.pptx')